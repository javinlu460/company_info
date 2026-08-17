# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

RED = RGBColor(0xB2, 0x2B, 0x2B)
DARK = RGBColor(0x1E, 0x1F, 0x20)
LIGHT = RGBColor(0xF2, 0xF3, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6A, 0x6D, 0x70)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_RED = RGBColor(0xFD, 0xF2, 0xF2)
BORDER_GRAY = RGBColor(0xE0, 0xE2, 0xDC)
LIGHT_TEXT = RGBColor(0xCC, 0xCC, 0xCC)
PALE_RED = RGBColor(0x99, 0x22, 0x22)
PALE_TEXT = RGBColor(0xFF, 0xDD, 0xDD)

def add_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size, font_color, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_slide(title, subtitle=None, bg_color=WHITE, title_color=DARK):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, bg_color)
    add_textbox(slide, 0.8, 0.6, 11.7, 1.0, title, 40, title_color, bold=True)
    if subtitle:
        add_textbox(slide, 0.8, 1.5, 11.7, 0.6, subtitle, 20, title_color if bg_color != RED else WHITE)
    slide_num = len(prs.slides)
    add_textbox(slide, 11.5, 0.3, 1.0, 0.4, f"{slide_num:02d} / 10", 12, 
                WHITE if bg_color in [RED, DARK] else GRAY, align=PP_ALIGN.RIGHT)
    return slide

def add_card(slide, left, top, width, height, title, desc, title_color=DARK, desc_color=GRAY, bg_color=WHITE, accent_color=RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if bg_color == WHITE:
        shape.line.color.rgb = BORDER_GRAY
    else:
        shape.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.08), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    add_textbox(slide, left + 0.25, top + 0.15, width - 0.4, 0.6, title, 18, title_color, bold=True)
    add_textbox(slide, left + 0.25, top + 0.75, width - 0.4, height - 1.0, desc, 13, desc_color)

# Slide 1: 封面
slide = add_slide("企业官网内容管理系统", 
                  "前后端分离 · 可配置化 · RBAC 权限",
                  bg_color=RED, title_color=WHITE)
add_textbox(slide, 0.8, 2.6, 11.7, 1.5, 
            "一套通用型企业官网后台 CMS，支持网站信息、产品、新闻、案例、FAQ、轮播图、留言等全站内容在线配置，零代码即可维护官网。",
            22, WHITE)

# Slide 2: 系统定位
slide = add_slide("系统定位", 
                  "面向企业官网的完整数字化管理平台",
                  bg_color=DARK, title_color=WHITE)

cards = [
    ("前台展示", "为访客提供企业介绍、产品服务、案例、新闻、FAQ、联系表单等完整官网浏览体验。"),
    ("后台配置", "为运营人员提供内容发布、栏目管理、网站配置、留言处理等一站式管理能力。"),
    ("权限管控", "基于 RBAC 模型的用户-角色-菜单权限控制，支持多角色协同与按钮级权限。")
]
for i, (title, desc) in enumerate(cards):
    add_card(slide, 0.8 + i * 4.2, 3.0, 3.9, 3.5, title, desc, 
             title_color=WHITE, desc_color=LIGHT_TEXT, bg_color=DARK_GRAY)

# Slide 3: 后台可配置模块总览
slide = add_slide("后台可配置模块总览", 
                  "所有前台展示内容均可通过后台在线管理",
                  bg_color=LIGHT, title_color=DARK)

modules = [
    ("网站配置", "公司信息、Logo、联系方式、备案号、富文本介绍等"),
    ("轮播图", "首页 Banner、标题、链接、排序、上下架"),
    ("产品管理", "产品分类、产品信息、封面图、详情、推荐、排序"),
    ("行业洞察", "文章分类、新闻/资讯内容、摘要、封面、状态"),
    ("解决方案", "方案标题、摘要、标签、详情、排序、状态"),
    ("客户案例", "案例标题、客户、行业、封面、描述、详情"),
    ("常见问题", "问题、答案、分类、排序、显示/隐藏"),
    ("留言管理", "访客留言、已读状态、回复内容")
]
for i, (title, desc) in enumerate(modules):
    row = i // 4
    col = i % 4
    left = 0.8 + col * 3.1
    top = 2.5 + row * 2.0
    add_card(slide, left, top, 2.9, 1.7, title, desc, title_color=DARK, desc_color=GRAY, bg_color=WHITE)

# Slide 4: 网站配置
slide = add_slide("网站配置", 
                  "统一维护企业品牌信息与联系资料",
                  bg_color=WHITE, title_color=DARK)

config_items = [
    ("企业名称", "公司名称、品牌简称，展示在网站头部与底部。"),
    ("企业 Logo", "上传网站 Logo，支持图片格式，自动应用到前台。"),
    ("联系方式", "电话、邮箱、地址、传真、微信号、工作时间、到厂路线。"),
    ("备案信息", "ICP 备案号等底部法律信息配置。"),
    ("公司介绍", "公司简介、业务与贸易范围、荣誉资质等富文本内容。"),
    ("SEO 信息", "网站标题、描述等基础搜索引擎优化配置。")
]
for i, (title, desc) in enumerate(config_items):
    left = 0.8 + (i % 2) * 6.0
    top = 2.4 + (i // 2) * 1.5
    add_card(slide, left, top, 5.7, 1.3, title, desc, title_color=DARK, desc_color=GRAY, bg_color=WHITE)

# Slide 5: 内容管理 - 产品与资讯
slide = add_slide("内容管理：产品与资讯", 
                  "结构化维护产品与行业内容",
                  bg_color=DARK, title_color=WHITE)

add_textbox(slide, 0.8, 2.4, 5.8, 0.6, "产品管理", 26, WHITE, bold=True)
product_items = [
    "支持产品分类增删改查",
    "产品名称、封面图、简介、详情、图集",
    "上下架状态控制",
    "推荐位设置与排序号",
    "浏览量统计"
]
for i, item in enumerate(product_items):
    add_textbox(slide, 1.0, 3.1 + i * 0.55, 5.6, 0.5, f"✓  {item}", 16, LIGHT_TEXT)

add_textbox(slide, 7.0, 2.4, 5.8, 0.6, "行业洞察 / 新闻", 26, WHITE, bold=True)
news_items = [
    "支持文章分类管理",
    "文章标题、摘要、封面图、富文本内容",
    "发布 / 草稿状态",
    "置顶、排序、作者、浏览量",
    "前台列表与详情展示"
]
for i, item in enumerate(news_items):
    add_textbox(slide, 7.2, 3.1 + i * 0.55, 5.6, 0.5, f"✓  {item}", 16, LIGHT_TEXT)

# Slide 6: 内容管理 - 解决方案、案例、FAQ
slide = add_slide("内容管理：方案 / 案例 / FAQ", 
                  "灵活配置业务展示与客户转化内容",
                  bg_color=LIGHT, title_color=DARK)

sections = [
    ("解决方案", [
        "方案标题、摘要、封面图",
        "标签 / 卖点配置",
        "详情富文本",
        "排序与上下架"
    ]),
    ("客户案例", [
        "案例标题、客户名称",
        "所属行业、封面图",
        "案例描述与详情",
        "排序与展示控制"
    ]),
    ("常见问题 FAQ", [
        "问题与答案编辑",
        "分类管理",
        "排序与显示/隐藏",
        "前台搜索与筛选"
    ])
]
for i, (title, items) in enumerate(sections):
    left = 0.8 + i * 4.2
    add_textbox(slide, left, 2.4, 3.9, 0.6, title, 24, DARK, bold=True)
    for j, item in enumerate(items):
        add_textbox(slide, left + 0.2, 3.1 + j * 0.6, 3.5, 0.5, f"✓  {item}", 15, GRAY)

# Slide 7: 用户与权限
slide = add_slide("用户与权限管理", 
                  "完整的 RBAC 权限体系，保障后台安全",
                  bg_color=RED, title_color=WHITE)

add_textbox(slide, 0.8, 2.4, 5.8, 0.6, "用户与角色", 26, WHITE, bold=True)
user_items = [
    "管理员账号管理",
    "角色创建与分配",
    "账号状态启用/禁用",
    "多角色协同支持"
]
for i, item in enumerate(user_items):
    add_textbox(slide, 1.0, 3.1 + i * 0.6, 5.6, 0.5, f"✓  {item}", 17, PALE_TEXT)

add_textbox(slide, 7.0, 2.4, 5.8, 0.6, "菜单与权限", 26, WHITE, bold=True)
perm_items = [
    "菜单目录与页面配置",
    "按钮级权限控制（增删改查）",
    "角色-菜单灵活授权",
    "权限实时生效",
    "操作审计留痕"
]
for i, item in enumerate(perm_items):
    add_textbox(slide, 7.2, 3.1 + i * 0.6, 5.6, 0.5, f"✓  {item}", 17, PALE_TEXT)

# Slide 8: 技术架构
slide = add_slide("技术架构", 
                  "成熟稳定的前后端分离架构，三端独立部署",
                  bg_color=LIGHT, title_color=DARK)

tech_stack = ["Vue 3", "Vite 5", "Element Plus", "Vue Router", "Spring Boot", "MyBatis-Plus", "MySQL", "RBAC"]
for i, tech in enumerate(tech_stack):
    left = 0.8 + (i % 4) * 3.0
    top = 2.6 + (i // 4) * 1.2
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.7), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER_GRAY
    add_textbox(slide, left, top + 0.25, 2.7, 0.5, tech, 16, DARK, align=PP_ALIGN.CENTER, bold=True)

shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_RED
shape.line.fill.background()
add_textbox(slide, 0.8, 5.7, 11.7, 0.8, 
            "三端分离：company-web（前台官网）+ company-cms（后台管理）+ company-admin（REST API）\n独立构建、独立部署、易于扩展",
            16, RED, align=PP_ALIGN.CENTER)

# Slide 9: 系统优势
slide = add_slide("系统优势", "", bg_color=DARK, title_color=WHITE)

advantages = [
    ("开箱即用", "提供完整数据库脚本与初始化数据，部署后即可使用。"),
    ("零代码配置", "网站信息、栏目、内容全部后台可视化配置，无需改代码。"),
    ("权限可控", "RBAC 多角色权限，保障后台操作安全。"),
    ("模块丰富", "覆盖企业官网常见内容类型，满足大部分展示需求。"),
    ("响应式前台", "适配桌面与移动端，保证多端访问体验。"),
    ("易于扩展", "模块化设计，可快速新增内容类型与页面。")
]
for i, (title, desc) in enumerate(advantages):
    left = 0.8 + (i % 3) * 4.2
    top = 2.0 + (i // 3) * 2.4
    add_card(slide, left, top, 3.9, 2.1, title, desc, 
             title_color=WHITE, desc_color=LIGHT_TEXT, bg_color=DARK_GRAY)

# Slide 10: 结尾
slide = add_slide("谢谢观看", 
                  "企业官网内容管理系统 · 让官网运营更简单",
                  bg_color=RED, title_color=WHITE)

add_textbox(slide, 0.8, 2.6, 11.7, 1.0, 
            "基于 Vue 3 + Spring Boot 的通用型企业官网 CMS，适用于制造、科技、服务等多种行业。",
            22, WHITE, align=PP_ALIGN.CENTER)

stats = [("8+", "可配模块"), ("7+", "内容类型"), ("3", "独立端"), ("1", "完整 RBAC")]
for i, (num, label) in enumerate(stats):
    left = 1.5 + i * 3.0
    add_textbox(slide, left, 4.0, 2.5, 1.0, num, 48, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left, 5.0, 2.5, 0.5, label, 16, PALE_TEXT, align=PP_ALIGN.CENTER)

prs.save('d:\\IdeaProjects\\company_info\\企业官网宣传介绍.pptx')
print("PPT 生成成功：d:\\IdeaProjects\\company_info\\企业官网宣传介绍.pptx")
